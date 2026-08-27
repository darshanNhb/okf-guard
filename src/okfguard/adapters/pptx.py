"""Adapter for PowerPoint presentations (.pptx) — speaker notes and off-canvas shapes.

PowerPoint has two distinct risk vectors for hidden content, different
from document-oriented formats:

1. **Speaker notes** — text in the notes pane is never seen by a live
   audience watching the presentation, but is commonly parsed by tools
   and generators building OKF bundles from decks.  Notes text is added
   to ``hidden_spans`` (not ``text``), since folding it into visible
   text would misrepresent what an actual audience would perceive.

2. **Off-canvas shapes** — shapes positioned entirely outside the
   visible slide boundary carry text that is invisible in any normal
   presentation view.  A shape is considered off-canvas if it has no
   visible overlap with the slide rectangle.

3. **White-text-on-fill colour matching** — for on-canvas, visible
   shapes, the text colour is compared against the shape's explicit fill
   colour.  Only explicit RGB values are compared; theme-indexed colours
   are skipped (false negative preferred over false positive), consistent
   with the principle that the pattern-matching layer provides a second
   line of defence.

**Requires:** ``python-pptx`` — install via
``pip install okf-guard[pptx]``.
"""

from __future__ import annotations

import io
import os

from okfguard.adapters.base import SourceAdapter, _utc_now_iso
from okfguard.core.models import ExtractedContent

# Colour tolerance per channel.
_COLOUR_TOLERANCE = 10


# ---------------------------------------------------------------------------
# Colour helpers
# ---------------------------------------------------------------------------

def _rgbcolor_to_tuple(rgb: object) -> tuple[int, int, int] | None:
    """Convert a python-pptx ``RGBColor`` to an ``(R, G, B)`` tuple.

    Returns ``None`` if *rgb* is ``None`` or cannot be converted.
    ``RGBColor`` is a subclass of ``int`` representing a 24-bit colour.
    """
    if rgb is None:
        return None
    try:
        val = int(rgb)  # type: ignore[arg-type,call-overload]
        return ((val >> 16) & 0xFF, (val >> 8) & 0xFF, val & 0xFF)
    except (TypeError, ValueError):
        return None


def _colours_close(
    c1: tuple[int, int, int],
    c2: tuple[int, int, int],
) -> bool:
    """Return ``True`` if two RGB colours are within tolerance."""
    return all(abs(a - b) <= _COLOUR_TOLERANCE for a, b in zip(c1, c2))


def _get_shape_fill_rgb(shape: object) -> tuple[int, int, int] | None:
    """Extract explicit RGB fill colour from a shape.

    Returns ``None`` for no fill, theme fill, or any value that cannot
    be confidently resolved to RGB.
    """
    try:
        fill = shape.fill  # type: ignore[attr-defined]
        # python-pptx fill.type: None (no fill), MSO_FILL_TYPE.SOLID,
        # MSO_FILL_TYPE.BACKGROUND, etc.
        if fill.type is None:
            return None
        # Attempt to access fore_color — may raise if fill is not solid
        fc = fill.fore_color
        if fc is None:
            return None
        # Only use explicit RGB, not theme colours
        # fc.type can be MSO_THEME_COLOR or RGB
        try:
            from pptx.enum.dml import MSO_THEME_COLOR
            if fc.type == MSO_THEME_COLOR:
                return None
        except (ImportError, AttributeError):
            pass
        return _rgbcolor_to_tuple(fc.rgb)
    except (AttributeError, TypeError, ValueError):
        return None


def _get_run_colour_rgb(run: object) -> tuple[int, int, int] | None:
    """Extract explicit RGB colour from a python-pptx text run.

    Returns ``None`` for theme colours, inherited colours, or when the
    colour property is not set.
    """
    try:
        font = run.font  # type: ignore[attr-defined]
        color_obj = font.color
        if color_obj is None:
            return None
        # Check if colour type is theme-based
        try:
            from pptx.enum.dml import MSO_THEME_COLOR
            if color_obj.type == MSO_THEME_COLOR:
                return None
        except (ImportError, AttributeError):
            pass
        return _rgbcolor_to_tuple(color_obj.rgb)
    except (AttributeError, TypeError, ValueError):
        return None


def _is_theme_colour(color_obj: object) -> bool:
    """Return ``True`` if *color_obj* uses a theme colour reference."""
    if color_obj is None:
        return False
    try:
        from pptx.enum.dml import MSO_THEME_COLOR
        return getattr(color_obj, "type", None) == MSO_THEME_COLOR
    except (ImportError, AttributeError):
        return False


# ---------------------------------------------------------------------------
# Off-canvas detection
# ---------------------------------------------------------------------------

def _is_off_canvas(
    shape: object,
    slide_width: int,
    slide_height: int,
) -> bool:
    """Return ``True`` if *shape* is positioned entirely outside the slide.

    A shape is off-canvas if it has no visible overlap with the slide
    rectangle — fully to the left, right, above, or below the visible
    area.  Coordinates are in EMUs (English Metric Units).
    """
    try:
        left = shape.left  # type: ignore[attr-defined]
        top = shape.top  # type: ignore[attr-defined]
        width = shape.width  # type: ignore[attr-defined]
        height = shape.height  # type: ignore[attr-defined]
    except (AttributeError, TypeError):
        return False

    if left is None or top is None or width is None or height is None:
        return False

    # Fully off one of the four edges
    if left + width <= 0:
        return True  # entirely to the left
    if top + height <= 0:
        return True  # entirely above
    if left >= slide_width:
        return True  # entirely to the right
    if top >= slide_height:
        return True  # entirely below
    return False


# ---------------------------------------------------------------------------
# PPTXAdapter
# ---------------------------------------------------------------------------

class PPTXAdapter(SourceAdapter):
    """Extract visible text and detect hidden content in PowerPoint files.

    **Requires:** ``python-pptx``.
    Install via ``pip install okf-guard[pptx]``.
    """

    @property
    def format_name(self) -> str:
        """Return ``'pptx'``."""
        return "pptx"

    def extract(self, source: str | bytes) -> ExtractedContent:
        """Extract content from a ``.pptx`` file.

        Args:
            source: A file path (``str``) or raw ``.pptx`` bytes.

        Returns:
            ``ExtractedContent`` with visible on-canvas text separated
            from speaker notes and off-canvas content, plus metadata
            including slide count.

        Raises:
            FileNotFoundError: If *source* is a path that doesn't exist.
            ValueError: If the content cannot be parsed as a valid
                ``.pptx`` file.
        """
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError(
                "PPTXAdapter requires 'python-pptx'.  "
                "Install it with:  pip install okf-guard[pptx]"
            ) from exc

        path: str | None = None

        try:
            if isinstance(source, bytes):
                prs = Presentation(io.BytesIO(source))
            elif isinstance(source, str):
                if not os.path.isfile(source):
                    raise FileNotFoundError(
                        f"PPTX file not found: {source!r}"
                    )
                path = source
                prs = Presentation(source)
            else:
                raise TypeError(
                    f"PPTXAdapter.extract() expects str or bytes, "
                    f"got {type(source).__name__}"
                )
        except (FileNotFoundError, TypeError, ImportError):
            raise
        except Exception as exc:
            raise ValueError(
                f"Failed to parse PPTX: {exc}"
            ) from exc

        slide_width = prs.slide_width or 0
        slide_height = prs.slide_height or 0
        slide_count = len(prs.slides)

        visible_parts: list[str] = []
        hidden_spans: list[str] = []

        # Auditability counter: how many colour comparisons were skipped
        # because at least one colour was theme-indexed and could not be
        # resolved to RGB.  This is an honest-bookkeeping signal (see
        # provenance design principle) — it does not affect risk scoring,
        # but lets a user see "N shapes had colours I couldn't check"
        # rather than silently assuming full coverage.
        # Deliberate decision to include this counter in v0.1.0 for
        # transparency, consistent with the project's principle of never
        # implying more certainty than actually exists.
        unresolved_theme_colors = 0

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_visible: list[str] = []

            # --- 1. Speaker notes → hidden_spans ---
            try:
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    notes_text = notes_frame.text
                    if notes_text and notes_text.strip():
                        hidden_spans.append(
                            f"[speaker notes, slide {slide_idx}] "
                            f"{notes_text}"
                        )
            except Exception:
                # Guard against notes-related errors in malformed files.
                pass

            # --- 2 & 3. Shapes: off-canvas and colour matching ---
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                shape_text = shape.text_frame.text
                if not shape_text or not shape_text.strip():
                    continue

                # Off-canvas check
                if _is_off_canvas(shape, slide_width, slide_height):
                    hidden_spans.append(
                        f"[slide {slide_idx}, off-canvas shape] "
                        f"{shape_text}"
                    )
                    continue

                # On-canvas — check for colour matching on runs
                shape_fill_rgb = _get_shape_fill_rgb(shape)
                shape_has_hidden_colour = False

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if not run.text.strip():
                            continue

                        run_rgb = _get_run_colour_rgb(run)

                        if run_rgb is not None and shape_fill_rgb is not None:
                            if _colours_close(run_rgb, shape_fill_rgb):
                                hidden_spans.append(
                                    f"[slide {slide_idx}, shape — "
                                    f"text color matches fill] {run.text}"
                                )
                                shape_has_hidden_colour = True
                                continue
                        elif (
                            _is_theme_colour(
                                getattr(run.font, "color", None)
                            )
                            or (
                                shape_fill_rgb is None
                                and _get_shape_fill_rgb(shape) is None
                                and _is_theme_fill(shape)
                            )
                        ):
                            unresolved_theme_colors += 1

                if not shape_has_hidden_colour:
                    slide_visible.append(shape_text)

            if slide_visible:
                # Separate slides with a blank line for readability.
                visible_parts.append(
                    f"--- Slide {slide_idx} ---\n"
                    + "\n".join(slide_visible)
                )

        metadata: dict[str, str] = {
            "format": self.format_name,
            "extracted_at": _utc_now_iso(),
            "slide_count": str(slide_count),
        }
        if path is not None:
            metadata["path"] = path
        if unresolved_theme_colors > 0:
            metadata["unresolved_theme_colors"] = str(
                unresolved_theme_colors
            )

        return ExtractedContent(
            text="\n\n".join(visible_parts),
            hidden_spans=hidden_spans,
            source_metadata=metadata,
        )


def _is_theme_fill(shape: object) -> bool:
    """Return ``True`` if the shape's fill is theme-based."""
    try:
        fill = shape.fill  # type: ignore[attr-defined]
        if fill.type is None:
            return False
        fc = fill.fore_color
        return _is_theme_colour(fc)
    except (AttributeError, TypeError):
        return False
